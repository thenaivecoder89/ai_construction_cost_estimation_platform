# File_name: extract_documents.py
# Purpose:
# Extract native text and table content from files listed in build_document_inventory.
#
# Output tables used:
# 1. documents
# 2. extracted_text
# 3. extracted_tables
# 4. extracted_table_rows
#
# Notes:
# - Bare-bones POC version.
# - No classes.
# - No schema creation.
# - No OCR. Scanned PDFs will be flagged as no_text_found_possible_scanned_pdf.

from pathlib import Path
import base64
import json
import os

import pandas as pd
from sqlalchemy import create_engine, text
from openai import OpenAI

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup

from global_rag.scripts import config


def clean_text(value):
    if value is None:
        return ""

    value = str(value).replace("\x00", " ")
    lines = [line.strip() for line in value.splitlines()]
    lines = [line for line in lines if line]

    return "\n".join(lines)


def estimate_tokens(value):
    value = clean_text(value)

    if value == "":
        return 0

    return max(1, int(len(value) / 4))


def clean_dataframe(df):
    df = df.copy()
    df = df.fillna("")

    new_columns = []
    seen_columns = {}

    for i, col in enumerate(df.columns, start=1):
        col_name = clean_text(col)

        if col_name == "":
            col_name = f"column_{i}"

        if col_name in seen_columns:
            seen_columns[col_name] += 1
            col_name = f"{col_name}_{seen_columns[col_name]}"
        else:
            seen_columns[col_name] = 1

        new_columns.append(col_name)

    df.columns = new_columns

    for col in df.columns:
        df[col] = df[col].astype(str)

    return df


def dataframe_to_jsonb_rows(table_id, document_id, sheet_name, df):
    output_rows = []
    df = clean_dataframe(df)

    for row_index, row in df.iterrows():
        row_dict = {}

        for col in df.columns:
            row_dict[str(col)] = clean_text(row[col])

        output_rows.append(
            {
                "table_id": table_id,
                "document_id": document_id,
                "sheet_name": sheet_name,
                "row_number": int(row_index) + 1,
                "row_data": json.dumps(row_dict, ensure_ascii=False),
            }
        )

    return output_rows


def extract_svg_text(file_path):
    svg_content = file_path.read_text(
        encoding="utf-8",
        errors="ignore"
    )

    soup = BeautifulSoup(svg_content, "xml")
    text_parts = []

    for tag_name in ["title", "desc", "text", "tspan"]:
        for tag in soup.find_all(tag_name):
            tag_text = clean_text(tag.get_text(separator=" "))
            if tag_text:
                text_parts.append(tag_text)

    if not text_parts:
        fallback_soup = BeautifulSoup(svg_content, "html.parser")
        text_parts.append(fallback_soup.get_text(separator="\n"))

    return clean_text("\n".join(text_parts))


def extract_dxf_text(file_path):
    """
    Extract visible text entities from ASCII DXF files using group-code pairs.
    This intentionally avoids CAD geometry parsing; it captures TEXT/MTEXT/
    ATTRIB strings that can be chunked and retrieved.
    """
    dxf_content = file_path.read_text(
        encoding="utf-8",
        errors="ignore"
    )

    lines = [line.rstrip("\n\r") for line in dxf_content.splitlines()]
    text_parts = []
    active_entity = None
    capture_next_value = False

    text_entity_types = {"TEXT", "MTEXT", "ATTRIB", "ATTDEF"}
    text_group_codes = {"1", "3"}

    for index in range(0, len(lines), 2):
        code = lines[index].strip()
        value = lines[index + 1].strip() if index + 1 < len(lines) else ""

        if code == "0":
            active_entity = value.upper()
            capture_next_value = False
            continue

        if active_entity in text_entity_types and code in text_group_codes:
            text_value = clean_text(value.replace("\\P", "\n"))
            if text_value:
                text_parts.append(text_value)
            capture_next_value = False
            continue

        if capture_next_value:
            text_value = clean_text(value.replace("\\P", "\n"))
            if text_value:
                text_parts.append(text_value)
            capture_next_value = False

    return clean_text("\n".join(text_parts))


def image_mime_type(file_extension):
    if file_extension in [".jpg", ".jpeg"]:
        return "image/jpeg"
    if file_extension == ".webp":
        return "image/webp"

    return "image/png"


def describe_image_file(file_path, file_extension, config_base):
    if not config.bool_env("EXTRACT_IMAGE_DESCRIPTIONS", True):
        return clean_text(
            f"Image file: {file_path.name}. "
            "Image description extraction disabled by EXTRACT_IMAGE_DESCRIPTIONS."
        )

    openai_api_key = config_base.get("openai_api_key")
    if not openai_api_key:
        raise RuntimeError("OPENAI_API_KEY is required to extract image descriptions.")

    image_model = os.getenv(
        "OPENAI_IMAGE_EXTRACTION_MODEL",
        config_base.get("llm_model", "gpt-4.1-mini"),
    )

    image_bytes = file_path.read_bytes()
    image_base64 = base64.b64encode(image_bytes).decode("ascii")
    mime_type = image_mime_type(file_extension)
    client = OpenAI(api_key=openai_api_key)

    response = client.responses.create(
        model=image_model,
        input=[
            {
                "role": "user",
                "content": [
                    {
                        "type": "input_text",
                        "text": (
                            "Extract the construction-cost-relevant information from this image. "
                            "Describe visible drawing labels, quantities, dimensions, scope notes, "
                            "cost assumptions, tables, legends, risks, and any text you can read. "
                            "Do not invent unreadable details."
                        ),
                    },
                    {
                        "type": "input_image",
                        "image_url": f"data:{mime_type};base64,{image_base64}",
                    },
                ],
            }
        ],
        max_output_tokens=int(os.getenv("OPENAI_IMAGE_EXTRACTION_MAX_OUTPUT_TOKENS", "1200")),
        store=False,
    )

    if hasattr(response, "output_text") and response.output_text:
        return clean_text(response.output_text)

    try:
        return clean_text(response.output[0].content[0].text)
    except Exception:
        return clean_text(str(response))


def sync_documents_from_inventory(engine):
    with engine.begin() as conn:
        conn.execute(
            text("""
                INSERT INTO documents (
                    document_id,
                    corpus_zone,
                    corpus_pack,
                    workstream,
                    source_folder,
                    relative_path,
                    file_name,
                    file_extension,
                    file_size_bytes,
                    file_checksum_sha256,
                    document_type,
                    document_title,
                    source_authority,
                    confidentiality_level,
                    is_client_confidential,
                    index_in_rag,
                    extraction_required,
                    ingest_status,
                    extraction_status,
                    extraction_quality,
                    error_message
                )
                SELECT
                    document_id,
                    source_group AS corpus_zone,
                    corpus_pack,
                    'ai_construction_cost_estimation_platform' AS workstream,
                    source_group AS source_folder,
                    relative_path,
                    file_name,
                    file_extension,
                    file_size_bytes,
                    sha256_checksum AS file_checksum_sha256,
                    file_extension AS document_type,
                    regexp_replace(file_name, '\\.[^.]*$', '') AS document_title,
                    corpus_pack AS source_authority,
                    CASE
                        WHEN source_group = 'client_data' THEN 'client_confidential'
                        ELSE 'reference_or_synthetic'
                    END AS confidentiality_level,
                    CASE
                        WHEN source_group = 'client_data' THEN TRUE
                        ELSE FALSE
                    END AS is_client_confidential,
                    CASE
                        WHEN index_in_rag = 'Yes' THEN TRUE
                        ELSE FALSE
                    END AS index_in_rag,
                    CASE
                        WHEN supported_file_type = 'Yes' THEN TRUE
                        ELSE FALSE
                    END AS extraction_required,
                    ingest_status,
                    'pending' AS extraction_status,
                    NULL AS extraction_quality,
                    NULL AS error_message
                FROM build_document_inventory
                WHERE supported_file_type = 'Yes'
                ON CONFLICT (document_id)
                DO UPDATE SET
                    corpus_zone = EXCLUDED.corpus_zone,
                    corpus_pack = EXCLUDED.corpus_pack,
                    workstream = EXCLUDED.workstream,
                    source_folder = EXCLUDED.source_folder,
                    relative_path = EXCLUDED.relative_path,
                    file_name = EXCLUDED.file_name,
                    file_extension = EXCLUDED.file_extension,
                    file_size_bytes = EXCLUDED.file_size_bytes,
                    file_checksum_sha256 = EXCLUDED.file_checksum_sha256,
                    document_type = EXCLUDED.document_type,
                    document_title = EXCLUDED.document_title,
                    source_authority = EXCLUDED.source_authority,
                    confidentiality_level = EXCLUDED.confidentiality_level,
                    is_client_confidential = EXCLUDED.is_client_confidential,
                    index_in_rag = EXCLUDED.index_in_rag,
                    extraction_required = EXCLUDED.extraction_required,
                    updated_at = NOW();
            """)
        )


def extract_documents(client_data: str, rebuild_inventory: str = "Y"):
    config_base = config.config_base()
    config_paths = config.config_paths(client_data=client_data)
    rebuild_inventory = rebuild_inventory.strip().upper()

    if rebuild_inventory not in ["Y", "N"]:
        raise ValueError("rebuild_inventory must be 'Y' or 'N'.")

    engine = create_engine(
        url=config_base["db_url"],
        pool_pre_ping=True
    )
    extract_pdf_tables = os.getenv("EXTRACT_PDF_TABLES", "N").strip().upper() == "Y"

    # Important: extracted_text / extracted_tables depend on documents.document_id
    sync_documents_from_inventory(engine)

    inventory_status_filter = ""
    if rebuild_inventory == "N":
        inventory_status_filter = "AND ingest_status = 'pending'"

    inventory_sql = f"""
        SELECT *
        FROM build_document_inventory
        WHERE supported_file_type = 'Yes'
          {inventory_status_filter}
        ORDER BY document_id;
    """

    inventory_df = pd.read_sql(text(inventory_sql), engine)

    documents_processed = 0
    documents_failed = 0
    extracted_text_rows_count = 0
    extracted_tables_count = 0
    extracted_table_rows_count = 0

    for _, inv_row in inventory_df.iterrows():

        document_id = str(inv_row["document_id"])
        file_extension = str(inv_row["file_extension"]).lower()
        relative_path = str(inv_row["relative_path"])
        absolute_path = str(inv_row["absolute_path"])

        file_path = Path(absolute_path)

        if not file_path.exists():
            fallback_path = config_paths["project_root"] / relative_path
            file_path = fallback_path

        text_rows = []
        table_rows = []
        table_data_rows = []

        table_counter = 1
        extraction_status = "extracted"
        document_extraction_quality = "extracted"
        error_message = None

        try:
            if not file_path.exists():
                extraction_status = "failed"
                document_extraction_quality = "failed_file_not_found"
                error_message = f"File not found: {file_path}"

                text_rows.append(
                    {
                        "document_id": document_id,
                        "page_no": None,
                        "section_heading": "file_missing",
                        "extraction_method": "file_path_check",
                        "extraction_quality": "failed_file_not_found",
                        "token_count_estimate": 0,
                        "text_content": "",
                    }
                )

            elif file_extension == ".pdf":

                pdf_doc = fitz.open(str(file_path))

                for page_no, page in enumerate(pdf_doc, start=1):
                    page_text = page.get_text("text")
                    cleaned_page_text = clean_text(page_text)

                    if cleaned_page_text == "":
                        extraction_quality = "no_text_found_possible_scanned_pdf"
                    else:
                        extraction_quality = "native_text_extracted"

                    text_rows.append(
                        {
                            "document_id": document_id,
                            "page_no": page_no,
                            "section_heading": f"page_{page_no}",
                            "extraction_method": "pymupdf_pdf_text",
                            "extraction_quality": extraction_quality,
                            "token_count_estimate": estimate_tokens(cleaned_page_text),
                            "text_content": cleaned_page_text,
                        }
                    )

                    if extract_pdf_tables:
                        try:
                            detected_tables = page.find_tables()

                            for detected_table in detected_tables.tables:
                                table_grid = detected_table.extract()

                                if table_grid and len(table_grid) > 0:
                                    df = pd.DataFrame(table_grid)
                                    df = clean_dataframe(df)

                                    table_id = f"{document_id}_TBL_{table_counter:04d}"

                                    table_rows.append(
                                        {
                                            "table_id": table_id,
                                            "document_id": document_id,
                                            "table_name": f"pdf_page_{page_no}_table_{table_counter}",
                                            "sheet_name": None,
                                            "page_no": page_no,
                                            "extracted_file_path": str(file_path),
                                            "row_count": len(df),
                                            "column_count": len(df.columns),
                                            "extraction_method": "pymupdf_find_tables",
                                            "extraction_quality": "table_extracted",
                                            "notes": "PDF table extracted using PyMuPDF page.find_tables().",
                                        }
                                    )

                                    table_data_rows.extend(
                                        dataframe_to_jsonb_rows(
                                            table_id=table_id,
                                            document_id=document_id,
                                            sheet_name=None,
                                            df=df,
                                        )
                                    )

                                    table_counter += 1

                        except Exception:
                            pass

                pdf_doc.close()

            elif file_extension == ".docx":

                doc = Document(str(file_path))
                doc_text_parts = []

                for section_no, section in enumerate(doc.sections, start=1):
                    for para in section.header.paragraphs:
                        para_text = clean_text(para.text)
                        if para_text:
                            doc_text_parts.append(f"[Header section {section_no}] {para_text}")

                    for para in section.footer.paragraphs:
                        para_text = clean_text(para.text)
                        if para_text:
                            doc_text_parts.append(f"[Footer section {section_no}] {para_text}")

                for para in doc.paragraphs:
                    para_text = clean_text(para.text)

                    if para_text:
                        style_name = ""

                        if para.style is not None:
                            style_name = str(para.style.name)

                        if style_name:
                            doc_text_parts.append(f"[{style_name}] {para_text}")
                        else:
                            doc_text_parts.append(para_text)

                text_content = "\n".join(doc_text_parts)

                text_rows.append(
                    {
                        "document_id": document_id,
                        "page_no": None,
                        "section_heading": "docx_document_text",
                        "extraction_method": "python_docx_text",
                        "extraction_quality": "native_text_extracted",
                        "token_count_estimate": estimate_tokens(text_content),
                        "text_content": clean_text(text_content),
                    }
                )

                for table_no, docx_table in enumerate(doc.tables, start=1):
                    table_grid = []

                    for table_row in docx_table.rows:
                        table_grid.append(
                            [clean_text(cell.text) for cell in table_row.cells]
                        )

                    if len(table_grid) > 0:
                        df = pd.DataFrame(table_grid)
                        df = clean_dataframe(df)

                        table_id = f"{document_id}_TBL_{table_counter:04d}"

                        table_rows.append(
                            {
                                "table_id": table_id,
                                "document_id": document_id,
                                "table_name": f"docx_table_{table_no}",
                                "sheet_name": None,
                                "page_no": None,
                                "extracted_file_path": str(file_path),
                                "row_count": len(df),
                                "column_count": len(df.columns),
                                "extraction_method": "python_docx_table",
                                "extraction_quality": "table_extracted",
                                "notes": "DOCX table extracted as raw cell grid.",
                            }
                        )

                        table_data_rows.extend(
                            dataframe_to_jsonb_rows(
                                table_id=table_id,
                                document_id=document_id,
                                sheet_name=None,
                                df=df,
                            )
                        )

                        table_counter += 1

            elif file_extension in [".pptx", ".ppt"]:

                presentation = Presentation(str(file_path))

                for slide_no, slide in enumerate(presentation.slides, start=1):
                    slide_text_parts = []

                    for shape in slide.shapes:

                        if hasattr(shape, "text"):
                            shape_text = clean_text(shape.text)

                            if shape_text:
                                slide_text_parts.append(shape_text)

                        if hasattr(shape, "has_table") and shape.has_table:
                            table_grid = []

                            for pptx_row in shape.table.rows:
                                table_grid.append(
                                    [clean_text(cell.text) for cell in pptx_row.cells]
                                )

                            if len(table_grid) > 0:
                                df = pd.DataFrame(table_grid)
                                df = clean_dataframe(df)

                                table_id = f"{document_id}_TBL_{table_counter:04d}"

                                table_rows.append(
                                    {
                                        "table_id": table_id,
                                        "document_id": document_id,
                                        "table_name": f"slide_{slide_no}_table_{table_counter}",
                                        "sheet_name": None,
                                        "page_no": slide_no,
                                        "extracted_file_path": str(file_path),
                                        "row_count": len(df),
                                        "column_count": len(df.columns),
                                        "extraction_method": "python_pptx_table",
                                        "extraction_quality": "table_extracted",
                                        "notes": "PPTX slide table extracted as raw cell grid.",
                                    }
                                )

                                table_data_rows.extend(
                                    dataframe_to_jsonb_rows(
                                        table_id=table_id,
                                        document_id=document_id,
                                        sheet_name=None,
                                        df=df,
                                    )
                                )

                                table_counter += 1

                    try:
                        if slide.has_notes_slide:
                            notes_text = clean_text(slide.notes_slide.notes_text_frame.text)

                            if notes_text:
                                slide_text_parts.append(f"[Speaker Notes]\n{notes_text}")
                    except Exception:
                        pass

                    slide_text = "\n".join(slide_text_parts)

                    text_rows.append(
                        {
                            "document_id": document_id,
                            "page_no": slide_no,
                            "section_heading": f"slide_{slide_no}",
                            "extraction_method": "python_pptx_text",
                            "extraction_quality": "slide_text_extracted",
                            "token_count_estimate": estimate_tokens(slide_text),
                            "text_content": clean_text(slide_text),
                        }
                    )

            elif file_extension in [".txt", ".md"]:

                text_content = file_path.read_text(
                    encoding="utf-8",
                    errors="ignore"
                )

                text_rows.append(
                    {
                        "document_id": document_id,
                        "page_no": None,
                        "section_heading": "text_file",
                        "extraction_method": "read_text_file",
                        "extraction_quality": "text_extracted",
                        "token_count_estimate": estimate_tokens(text_content),
                        "text_content": clean_text(text_content),
                    }
                )

            elif file_extension == ".html":

                html_content = file_path.read_text(
                    encoding="utf-8",
                    errors="ignore"
                )

                soup = BeautifulSoup(html_content, "html.parser")

                for tag in soup(["script", "style"]):
                    tag.decompose()

                html_text = soup.get_text(separator="\n")

                text_rows.append(
                    {
                        "document_id": document_id,
                        "page_no": None,
                        "section_heading": "html_document",
                        "extraction_method": "beautifulsoup_html_text",
                        "extraction_quality": "html_text_extracted",
                        "token_count_estimate": estimate_tokens(html_text),
                        "text_content": clean_text(html_text),
                    }
                )

                try:
                    html_tables = pd.read_html(str(file_path))

                    for html_table_no, df in enumerate(html_tables, start=1):
                        df = clean_dataframe(df)

                        table_id = f"{document_id}_TBL_{table_counter:04d}"

                        table_rows.append(
                            {
                                "table_id": table_id,
                                "document_id": document_id,
                                "table_name": f"html_table_{html_table_no}",
                                "sheet_name": None,
                                "page_no": None,
                                "extracted_file_path": str(file_path),
                                "row_count": len(df),
                                "column_count": len(df.columns),
                                "extraction_method": "pandas_read_html",
                                "extraction_quality": "table_extracted",
                                "notes": "HTML table extracted using pandas.read_html.",
                            }
                        )

                        table_data_rows.extend(
                            dataframe_to_jsonb_rows(
                                table_id=table_id,
                                document_id=document_id,
                                sheet_name=None,
                                df=df,
                            )
                        )

                        table_counter += 1

                except Exception:
                    pass

            elif file_extension == ".json":

                raw_json = file_path.read_text(
                    encoding="utf-8",
                    errors="ignore"
                )

                parsed_json = json.loads(raw_json)

                json_text = json.dumps(
                    parsed_json,
                    indent=2,
                    ensure_ascii=False,
                    default=str
                )

                text_rows.append(
                    {
                        "document_id": document_id,
                        "page_no": None,
                        "section_heading": "json_document",
                        "extraction_method": "json_loads",
                        "extraction_quality": "json_text_extracted",
                        "token_count_estimate": estimate_tokens(json_text),
                        "text_content": clean_text(json_text),
                    }
                )

                if isinstance(parsed_json, list) and len(parsed_json) > 0:
                    if all(isinstance(item, dict) for item in parsed_json):
                        df = pd.DataFrame(parsed_json)
                        df = clean_dataframe(df)

                        table_id = f"{document_id}_TBL_{table_counter:04d}"

                        table_rows.append(
                            {
                                "table_id": table_id,
                                "document_id": document_id,
                                "table_name": "json_root_records",
                                "sheet_name": None,
                                "page_no": None,
                                "extracted_file_path": str(file_path),
                                "row_count": len(df),
                                "column_count": len(df.columns),
                                "extraction_method": "json_list_to_table",
                                "extraction_quality": "table_extracted",
                                "notes": "Top-level JSON list converted to table rows.",
                            }
                        )

                        table_data_rows.extend(
                            dataframe_to_jsonb_rows(
                                table_id=table_id,
                                document_id=document_id,
                                sheet_name=None,
                                df=df,
                            )
                        )

                        table_counter += 1

                if isinstance(parsed_json, dict):
                    for json_key, json_value in parsed_json.items():
                        if isinstance(json_value, list) and len(json_value) > 0:
                            if all(isinstance(item, dict) for item in json_value):
                                df = pd.DataFrame(json_value)
                                df = clean_dataframe(df)

                                table_id = f"{document_id}_TBL_{table_counter:04d}"

                                table_rows.append(
                                    {
                                        "table_id": table_id,
                                        "document_id": document_id,
                                        "table_name": f"json_{json_key}",
                                        "sheet_name": None,
                                        "page_no": None,
                                        "extracted_file_path": str(file_path),
                                        "row_count": len(df),
                                        "column_count": len(df.columns),
                                        "extraction_method": "json_nested_list_to_table",
                                        "extraction_quality": "table_extracted",
                                        "notes": f"Nested JSON list converted to table rows from key: {json_key}.",
                                    }
                                )

                                table_data_rows.extend(
                                    dataframe_to_jsonb_rows(
                                        table_id=table_id,
                                        document_id=document_id,
                                        sheet_name=None,
                                        df=df,
                                    )
                                )

                                table_counter += 1

            elif file_extension == ".svg":

                svg_text = extract_svg_text(file_path)

                text_rows.append(
                    {
                        "document_id": document_id,
                        "page_no": None,
                        "section_heading": "svg_text",
                        "extraction_method": "beautifulsoup_svg_text",
                        "extraction_quality": "svg_text_extracted" if svg_text else "no_extractable_svg_text",
                        "token_count_estimate": estimate_tokens(svg_text),
                        "text_content": clean_text(svg_text),
                    }
                )

            elif file_extension == ".dxf":

                dxf_text = extract_dxf_text(file_path)

                text_rows.append(
                    {
                        "document_id": document_id,
                        "page_no": None,
                        "section_heading": "dxf_text_entities",
                        "extraction_method": "dxf_text_entity_parse",
                        "extraction_quality": "dxf_text_extracted" if dxf_text else "no_extractable_dxf_text",
                        "token_count_estimate": estimate_tokens(dxf_text),
                        "text_content": clean_text(dxf_text),
                    }
                )

            elif file_extension in config_paths["supported_image_extensions"]:

                image_description = describe_image_file(
                    file_path=file_path,
                    file_extension=file_extension,
                    config_base=config_base
                )

                text_rows.append(
                    {
                        "document_id": document_id,
                        "page_no": None,
                        "section_heading": "image_description",
                        "extraction_method": "openai_image_description",
                        "extraction_quality": "image_description_extracted",
                        "token_count_estimate": estimate_tokens(image_description),
                        "text_content": clean_text(image_description),
                    }
                )

            elif file_extension == ".csv":

                df = pd.read_csv(file_path)
                df = clean_dataframe(df)

                table_id = f"{document_id}_TBL_{table_counter:04d}"

                table_rows.append(
                    {
                        "table_id": table_id,
                        "document_id": document_id,
                        "table_name": file_path.stem,
                        "sheet_name": None,
                        "page_no": None,
                        "extracted_file_path": str(file_path),
                        "row_count": len(df),
                        "column_count": len(df.columns),
                        "extraction_method": "pandas_read_csv",
                        "extraction_quality": "table_extracted",
                        "notes": "CSV extracted into extracted_tables and extracted_table_rows.",
                    }
                )

                table_data_rows.extend(
                    dataframe_to_jsonb_rows(
                        table_id=table_id,
                        document_id=document_id,
                        sheet_name=None,
                        df=df,
                    )
                )

                table_counter += 1

            elif file_extension in [".xlsx", ".xls"]:

                excel_sheets = pd.read_excel(file_path, sheet_name=None)

                for sheet_name, df in excel_sheets.items():
                    df = clean_dataframe(df)

                    table_id = f"{document_id}_TBL_{table_counter:04d}"

                    table_rows.append(
                        {
                            "table_id": table_id,
                            "document_id": document_id,
                            "table_name": f"{file_path.stem}_{sheet_name}",
                            "sheet_name": str(sheet_name),
                            "page_no": None,
                            "extracted_file_path": str(file_path),
                            "row_count": len(df),
                            "column_count": len(df.columns),
                            "extraction_method": "pandas_read_excel",
                            "extraction_quality": "table_extracted",
                            "notes": "Excel sheet extracted into extracted_tables and extracted_table_rows.",
                        }
                    )

                    table_data_rows.extend(
                        dataframe_to_jsonb_rows(
                            table_id=table_id,
                            document_id=document_id,
                            sheet_name=str(sheet_name),
                            df=df,
                        )
                    )

                    table_counter += 1

            else:
                extraction_status = "skipped"
                document_extraction_quality = "unsupported_extension"
                error_message = f"Unsupported extension: {file_extension}"

            total_text_chars = sum(
                len(clean_text(row["text_content"]))
                for row in text_rows
            )

            total_table_data_rows = len(table_data_rows)

            if extraction_status != "failed" and extraction_status != "skipped":
                if total_text_chars == 0 and total_table_data_rows == 0:
                    extraction_status = "extracted_no_content"
                    document_extraction_quality = "no_extractable_content"
                else:
                    extraction_status = "extracted"
                    document_extraction_quality = "extracted"

        except Exception as e:
            extraction_status = "failed"
            document_extraction_quality = "failed"
            error_message = str(e)[:1000]
            documents_failed += 1

            text_rows = [
                {
                    "document_id": document_id,
                    "page_no": None,
                    "section_heading": "extraction_error",
                    "extraction_method": "extract_documents",
                    "extraction_quality": "failed",
                    "token_count_estimate": 0,
                    "text_content": clean_text(str(e))[:1000],
                }
            ]

            table_rows = []
            table_data_rows = []

        with engine.begin() as conn:

            conn.execute(
                text("""
                    DELETE FROM extracted_table_rows
                    WHERE document_id = :document_id
                """),
                {"document_id": document_id}
            )

            conn.execute(
                text("""
                    DELETE FROM extracted_tables
                    WHERE document_id = :document_id
                """),
                {"document_id": document_id}
            )

            conn.execute(
                text("""
                    DELETE FROM extracted_text
                    WHERE document_id = :document_id
                """),
                {"document_id": document_id}
            )

            if len(text_rows) > 0:
                conn.execute(
                    text("""
                        INSERT INTO extracted_text (
                            document_id,
                            page_no,
                            section_heading,
                            extraction_method,
                            extraction_quality,
                            token_count_estimate,
                            text_content
                        )
                        VALUES (
                            :document_id,
                            :page_no,
                            :section_heading,
                            :extraction_method,
                            :extraction_quality,
                            :token_count_estimate,
                            :text_content
                        )
                    """),
                    text_rows
                )

            if len(table_rows) > 0:
                conn.execute(
                    text("""
                        INSERT INTO extracted_tables (
                            table_id,
                            document_id,
                            table_name,
                            sheet_name,
                            page_no,
                            extracted_file_path,
                            row_count,
                            column_count,
                            extraction_method,
                            extraction_quality,
                            notes
                        )
                        VALUES (
                            :table_id,
                            :document_id,
                            :table_name,
                            :sheet_name,
                            :page_no,
                            :extracted_file_path,
                            :row_count,
                            :column_count,
                            :extraction_method,
                            :extraction_quality,
                            :notes
                        )
                    """),
                    table_rows
                )

            if len(table_data_rows) > 0:
                conn.execute(
                    text("""
                        INSERT INTO extracted_table_rows (
                            table_id,
                            document_id,
                            sheet_name,
                            row_number,
                            row_data
                        )
                        VALUES (
                            :table_id,
                            :document_id,
                            :sheet_name,
                            :row_number,
                            CAST(:row_data AS JSONB)
                        )
                    """),
                    table_data_rows
                )

            conn.execute(
                text("""
                    UPDATE documents
                    SET ingest_status = :ingest_status,
                        extraction_status = :extraction_status,
                        extraction_quality = :extraction_quality,
                        error_message = :error_message,
                        updated_at = NOW()
                    WHERE document_id = :document_id
                """),
                {
                    "document_id": document_id,
                    "ingest_status": extraction_status,
                    "extraction_status": extraction_status,
                    "extraction_quality": document_extraction_quality,
                    "error_message": error_message,
                }
            )

            conn.execute(
                text("""
                    UPDATE build_document_inventory
                    SET ingest_status = :ingest_status
                    WHERE document_id = :document_id
                """),
                {
                    "document_id": document_id,
                    "ingest_status": extraction_status,
                }
            )

        documents_processed += 1
        extracted_text_rows_count += len(text_rows)
        extracted_tables_count += len(table_rows)
        extracted_table_rows_count += len(table_data_rows)

    return {
        "message": "Document extraction completed.",
        "client_data": client_data,
        "mode": "rebuild" if rebuild_inventory == "Y" else "update",
        "documents_processed": documents_processed,
        "documents_failed": documents_failed,
        "extracted_text_rows": extracted_text_rows_count,
        "extracted_tables": extracted_tables_count,
        "extracted_table_rows": extracted_table_rows_count,
    }


if __name__ == "__main__":
    print(extract_documents(client_data="TXN_ADDC_001", rebuild_inventory="Y"))
